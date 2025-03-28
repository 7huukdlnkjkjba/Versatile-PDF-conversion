import tkinter as tk
from tkinter import ttk, filedialog, messagebox
from threading import Thread
from PIL import Image
from pdf2docx import Converter
from pdf2image import convert_from_path
import pytesseract
import PyPDF2
import pandas as pd
from pptx import Presentation
from docx import Document
import os
import sys


class FileConverterApp:
    def __init__(self, root):
        self.root = root
        self.root.title("全能文件转换器 v2.0")
        self.setup_ui()
        self.conversion_functions = {
            "image2pdf": self.image_to_pdf,
            "pdf2ppt": self.pdf_to_ppt,
            "ppt2word": self.ppt_to_word,
            "pdf2excel": self.pdf_to_excel,
            "word2excel": self.word_to_excel,
            "pdf2txt": self.pdf_to_txt
        }

    def setup_ui(self):
        main_frame = ttk.Frame(self.root, padding=20)
        main_frame.grid(row=0, column=0, sticky="nsew")

        # 文件选择
        ttk.Label(main_frame, text="源文件:").grid(row=0, column=0, sticky="w")
        self.source_entry = ttk.Entry(main_frame, width=40)
        self.source_entry.grid(row=0, column=1)
        ttk.Button(main_frame, text="浏览...", command=self.select_source).grid(row=0, column=2)

        # 转换类型
        ttk.Label(main_frame, text="转换类型:").grid(row=1, column=0, sticky="w")
        self.conversion_type = ttk.Combobox(main_frame, values=[
            "图片转PDF", "PDF转PPT", "PPT转Word",
            "PDF转Excel", "Word转Excel", "PDF转TXT"
        ])
        self.conversion_type.grid(row=1, column=1, sticky="ew")
        self.conversion_type.current(0)

        # 进度条
        self.progress = ttk.Progressbar(main_frame, mode="determinate")
        self.progress.grid(row=2, column=0, columnspan=3, sticky="ew", pady=10)

        # 操作按钮
        ttk.Button(main_frame, text="开始转换", command=self.start_conversion).grid(row=3, column=1, pady=10)

    def select_source(self):
        filetypes = [
            ("All files", "*.*"),
            ("PDF files", "*.pdf"),
            ("Image files", "*.jpg *.png *.jpeg"),
            ("PPT files", "*.pptx"),
            ("Word files", "*.docx")
        ]
        filename = filedialog.askopenfilename(filetypes=filetypes)
        self.source_entry.delete(0, tk.END)
        self.source_entry.insert(0, filename)

    def start_conversion(self):
        src_file = self.source_entry.get()
        if not src_file:
            messagebox.showerror("错误", "请先选择源文件！")
            return

        conv_type = self.conversion_type.get().replace(" ", "").lower()
        Thread(target=self.run_conversion, args=(src_file, conv_type)).start()

    def run_conversion(self, src, conv_type):
        self.progress["value"] = 0
        try:
            output_path = filedialog.asksaveasfilename(
                defaultextension=self.get_extension(conv_type),
                filetypes=[(f"{conv_type.upper()} files", self.get_extension(conv_type))]
            )
            if output_path:
                self.conversion_functions[conv_type](src, output_path)
                messagebox.showinfo("成功", "转换完成！")
        except Exception as e:
            messagebox.showerror("错误", str(e))
        finally:
            self.progress["value"] = 100

    # 各转换功能实现
    def image_to_pdf(self, src, dst):
        img = Image.open(src)
        img.save(dst, "PDF", resolution=100.0)

    def pdf_to_ppt(self, src, dst, dpi=200):
        images = convert_from_path(src, dpi=dpi)
        prs = Presentation()
        for img in images:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            img.save("temp_img.png")
            slide.shapes.add_picture("temp_img.png", 0, 0, prs.slide_width, prs.slide_height)
        prs.save(dst)
        os.remove("temp_img.png")

    def ppt_to_word(self, src, dst):
        prs = Presentation(src)
        doc = Document()
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    doc.add_paragraph(shape.text)
        doc.save(dst)

    def pdf_to_excel(self, src, dst):
        text = ""
        with open(src, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                text += page.extract_text()
        df = pd.DataFrame([line.split("\t") for line in text.split("\n")])
        df.to_excel(dst, index=False)

    def word_to_excel(self, src, dst):
        doc = Document(src)
        data = []
        for table in doc.tables:
            for row in table.rows:
                data.append([cell.text for cell in row.cells])
        pd.DataFrame(data).to_excel(dst, index=False)

    def pdf_to_txt(self, src, dst):
        with open(src, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            with open(dst, "w", encoding="utf-8") as out:
                for page in reader.pages:
                    out.write(page.extract_text())

    def get_extension(self, conv_type):
        return {
            "image2pdf": ".pdf",
            "pdf2ppt": ".pptx",
            "ppt2word": ".docx",
            "pdf2excel": ".xlsx",
            "word2excel": ".xlsx",
            "pdf2txt": ".txt"
        }[conv_type]


if __name__ == "__main__":
    root = tk.Tk()
    app = FileConverterApp(root)
    root.mainloop()